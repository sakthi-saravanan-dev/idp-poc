# from flask import Flask, request, render_template, send_file
# import os
# import fitz  # PyMuPDF
# import pandas as pd
# import tempfile
# import io
# from docx import Document as DocxDocument
# from werkzeug.utils import secure_filename
# import pptx
# from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE
# from collections import defaultdict
# from transformers import pipeline  # For summarization

# app = Flask(__name__)
# extracted_data = []

# # Local summarizer model (small for demo; replace with larger one if needed)
# summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

# def classify_block(text):
#     text = text.strip()
#     if not text:
#         return "Empty"
#     elif len(text.split()) < 6 and text.isupper():
#         return "Title"
#     elif "|" in text or text.count("\t") > 1:
#         return "Table"
#     elif len(text.split()) > 40:
#         return "Paragraph"
#     elif len(text) < 40:
#         return "String"
#     else:
#         return "Paragraph"

# def summarize_if_needed(text):
#     if len(text.split()) > 100:
#         summary = summarizer(text, max_length=100, min_length=30, do_sample=False)[0]['summary_text']
#         return summary
#     return text

# def extract_from_pdf(file, filename):
#     results = defaultdict(list)
#     with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
#         tmp.write(file.read())
#         tmp.flush()
#         doc = fitz.open(tmp.name)
#         for page_number, page in enumerate(doc, start=1):
#             page_key = f"{filename}::Page {page_number}"
#             blocks = page.get_text("blocks")
#             for block in blocks:
#                 text = block[4].strip()
#                 if text:
#                     content_type = classify_block(text)
#                     summarized = summarize_if_needed(text)
#                     results[page_key].append((content_type, summarized))
#             images = page.get_images(full=True)
#             for i, _ in enumerate(images):
#                 results[page_key].append(("Image", f"Image on page {page_number}, index {i}"))
#     return results

# def extract_from_docx(file, filename):
#     results = defaultdict(list)
#     doc = DocxDocument(file)
#     page_number = 1
#     page_key = f"{filename}::Page {page_number}"
#     count = 0
#     for para in doc.paragraphs:
#         text = para.text.strip()
#         if text:
#             count += 1
#             content_type = classify_block(text)
#             summarized = summarize_if_needed(text)
#             results[page_key].append((content_type, summarized))
#             if count % 10 == 0:  # simulate page
#                 page_number += 1
#                 page_key = f"{filename}::Page {page_number}"
#     return results

# def extract_from_csv(file, filename):
#     results = defaultdict(list)
#     df = pd.read_csv(file)
#     page_key = f"{filename}::Page 1"
#     table_text = df.to_csv(sep="\t", index=False)
#     results[page_key].append(("Table", table_text))
#     return results

# def extract_from_excel(file, filename):
#     results = defaultdict(list)
#     xls = pd.ExcelFile(file)
#     for sheet_name in xls.sheet_names:
#         df = xls.parse(sheet_name)
#         page_key = f"{filename}::{sheet_name}"
#         table_text = df.to_csv(sep="\t", index=False)
#         results[page_key].append(("Table", table_text))
#     return results

# def extract_from_txt(file, filename):
#     results = defaultdict(list)
#     content = file.read().decode('utf-8')
#     for i, line in enumerate(content.split('\n'), start=1):
#         line = line.strip()
#         if line:
#             page_key = f"{filename}::Line {i}"
#             content_type = classify_block(line)
#             summarized = summarize_if_needed(line)
#             results[page_key].append((content_type, summarized))
#     return results

# def extract_from_pptx(file, filename):
#     results = defaultdict(list)
#     prs = Presentation(file)
#     for i, slide in enumerate(prs.slides, start=1):
#         page_key = f"{filename}::Slide {i}"
#         for shape in slide.shapes:
#             if not shape.has_text_frame:
#                 continue
#             text = shape.text.strip()
#             if text:
#                 content_type = classify_block(text)
#                 summarized = summarize_if_needed(text)
#                 results[page_key].append((content_type, summarized))
#     return results

# @app.route('/', methods=['GET', 'POST'])
# def upload_files():
#     global extracted_data
#     if request.method == 'POST':
#         files = request.files.getlist('documents')
#         if not files or files[0].filename == '':
#             return render_template('index.html', error="No files selected")

#         extracted_data = []
#         for file in files:
#             filename = secure_filename(file.filename)
#             ext = os.path.splitext(filename)[1].lower()
#             page_data = {}

#             if ext == ".pdf":
#                 page_data = extract_from_pdf(file, filename)
#             elif ext == ".docx":
#                 page_data = extract_from_docx(file, filename)
#             elif ext == ".csv":
#                 page_data = extract_from_csv(file, filename)
#             elif ext in [".xls", ".xlsx"]:
#                 page_data = extract_from_excel(file, filename)
#             elif ext == ".txt":
#                 page_data = extract_from_txt(file, filename)
#             elif ext == ".pptx":
#                 page_data = extract_from_pptx(file, filename)
#             else:
#                 page_data[f"{filename}::Page 1"] = [("Unsupported", f"File type {ext} is not supported")]

#             for key, content in page_data.items():
#                 fname, page = key.split("::")
#                 extracted_data.append({
#                     "Filename": fname,
#                     "Page Number": page,
#                     "Content": content
#                 })

#     return render_template('index.html', data=extracted_data)

# @app.route('/download')
# def download_docx():
#     from docx import Document
#     doc = Document()

#     for row in extracted_data:
#         doc.add_heading(f"{row['Filename']} - {row['Page Number']}", level=2)
#         for ctype, content in row["Content"]:
#             if ctype == "Title":
#                 doc.add_heading(content, level=3)
#             elif ctype == "Paragraph":
#                 doc.add_paragraph(content)
#             elif ctype == "Table":
#                 rows = content.strip().split('\n')
#                 if rows:
#                     cols = rows[0].split('\t')
#                     table = doc.add_table(rows=1, cols=len(cols))
#                     hdr_cells = table.rows[0].cells
#                     for i, col in enumerate(cols):
#                         hdr_cells[i].text = col
#                     for row_text in rows[1:]:
#                         row_data = row_text.split('\t')
#                         row_cells = table.add_row().cells
#                         for i, item in enumerate(row_data):
#                             row_cells[i].text = item
#             else:
#                 doc.add_paragraph(f"[{ctype}] {content}", style='Intense Quote')
#         doc.add_paragraph("\n")

#     output = io.BytesIO()
#     doc.save(output)
#     output.seek(0)
#     return send_file(output, as_attachment=True, download_name="extracted_summary.docx")

# if __name__ == '__main__':
#     app.run(debug=True)

from flask import Flask, request, render_template, send_file
import os
import fitz  # PyMuPDF
import pandas as pd
import tempfile
import io
from docx import Document as DocxDocument
from werkzeug.utils import secure_filename
import pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import defaultdict
from transformers import pipeline

app = Flask(__name__)
extracted_data = []

# Summarizer
summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")

def classify_block(text):
    text = text.strip()
    if not text:
        return "Empty"
    elif len(text.split()) < 6 and text.isupper():
        return "Title"
    elif "|" in text or text.count("\t") > 1:
        return "Table"
    elif len(text.split()) > 40:
        return "Paragraph"
    elif len(text) < 40:
        return "String"
    else:
        return "Paragraph"

def summarize_if_needed(text):
    if len(text.split()) > 100:
        summary = summarizer(text, max_length=100, min_length=30, do_sample=False)[0]['summary_text']
        return summary
    return text

def extract_from_pdf(file, filename):
    results = defaultdict(list)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file.read())
        tmp.flush()
        doc = fitz.open(tmp.name)
        for page_number, page in enumerate(doc, start=1):
            page_key = f"{filename}::Page {page_number}"
            # Sort blocks: (y0, x0)
            blocks = sorted(page.get_text("blocks"), key=lambda b: (b[1], b[0]))
            for block in blocks:
                text = block[4].strip()
                if text:
                    content_type = classify_block(text)
                    summarized = summarize_if_needed(text)
                    results[page_key].append((content_type, summarized))
            images = page.get_images(full=True)
            for i, _ in enumerate(images):
                results[page_key].append(("Image", f"Image on page {page_number}, index {i}"))
    return results

def extract_from_docx(file, filename):
    results = defaultdict(list)
    doc = DocxDocument(file)
    page_number = 1
    page_key = f"{filename}::Page {page_number}"
    count = 0
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            count += 1
            content_type = classify_block(text)
            summarized = summarize_if_needed(text)
            results[page_key].append((content_type, summarized))
            if count % 10 == 0:
                page_number += 1
                page_key = f"{filename}::Page {page_number}"
    return results

def extract_from_csv(file, filename):
    results = defaultdict(list)
    df = pd.read_csv(file)
    page_key = f"{filename}::Page 1"
    table_text = df.to_csv(sep="\t", index=False)
    results[page_key].append(("Table", table_text))
    return results

def extract_from_excel(file, filename):
    results = defaultdict(list)
    xls = pd.ExcelFile(file)
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        page_key = f"{filename}::{sheet_name}"
        table_text = df.to_csv(sep="\t", index=False)
        results[page_key].append(("Table", table_text))
    return results

def extract_from_txt(file, filename):
    results = defaultdict(list)
    content = file.read().decode('utf-8')
    for i, line in enumerate(content.split('\n'), start=1):
        line = line.strip()
        if line:
            page_key = f"{filename}::Line {i}"
            content_type = classify_block(line)
            summarized = summarize_if_needed(line)
            results[page_key].append((content_type, summarized))
    return results

def extract_from_pptx(file, filename):
    results = defaultdict(list)
    prs = Presentation(file)
    for i, slide in enumerate(prs.slides, start=1):
        page_key = f"{filename}::Slide {i}"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            if text:
                content_type = classify_block(text)
                summarized = summarize_if_needed(text)
                results[page_key].append((content_type, summarized))
    return results

@app.route('/', methods=['GET', 'POST'])
def upload_files():
    global extracted_data

    if request.method == 'POST':
        files = request.files.getlist('documents')
        if not files or files[0].filename == '':
            return render_template('index.html', error="No files selected")

        extracted_data = []
        for file in files:
            filename = secure_filename(file.filename)
            ext = os.path.splitext(filename)[1].lower()
            page_data = {}

            if ext == ".pdf":
                page_data = extract_from_pdf(file, filename)
            elif ext == ".docx":
                page_data = extract_from_docx(file, filename)
            elif ext == ".csv":
                page_data = extract_from_csv(file, filename)
            elif ext in [".xls", ".xlsx"]:
                page_data = extract_from_excel(file, filename)
            elif ext == ".txt":
                page_data = extract_from_txt(file, filename)
            elif ext == ".pptx":
                page_data = extract_from_pptx(file, filename)
            else:
                page_data[f"{filename}::Page 1"] = [("Unsupported", f"File type {ext} is not supported")]

            for key, content in page_data.items():
                fname, page = key.split("::")
                extracted_data.append({
                    "Filename": fname,
                    "Page Number": page,
                    "Content": content
                })

        return render_template('index.html', data=extracted_data)

    # Reset data on page load
    extracted_data = []
    return render_template('index.html')

@app.route('/reset', methods=['GET'])
def reset():
    global extracted_data
    extracted_data = []
    return render_template('index.html')

@app.route('/download')
def download_docx():
    from docx import Document
    doc = Document()

    for row in extracted_data:
        doc.add_heading(f"{row['Filename']} - {row['Page Number']}", level=2)
        for ctype, content in row["Content"]:
            if ctype == "Title":
                doc.add_heading(content, level=3)
            elif ctype == "Paragraph":
                doc.add_paragraph(content)
            elif ctype == "Table":
                rows = content.strip().split('\n')
                if rows:
                    cols = rows[0].split('\t')
                    table = doc.add_table(rows=1, cols=len(cols))
                    hdr_cells = table.rows[0].cells
                    for i, col in enumerate(cols):
                        hdr_cells[i].text = col
                    for row_text in rows[1:]:
                        row_data = row_text.split('\t')
                        row_cells = table.add_row().cells
                        for i, item in enumerate(row_data):
                            row_cells[i].text = item
            else:
                doc.add_paragraph(f"[{ctype}] {content}", style='Intense Quote')
        doc.add_paragraph("\n")

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return send_file(output, as_attachment=True, download_name="extracted_summary.docx")

if __name__ == '__main__':
    app.run(debug=True)